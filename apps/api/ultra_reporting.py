"""
ULTRA REPORTING MODULE
Kombinim i Excel, PowerPoint, Datadog, Grafana, Prometheus
Gjeneron raporte sipas kërkesës - Metrics -> Excel -> PowerPoint -> Dashboards
"""

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.chart import LineChart, BarChart, PieChart, Reference
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta
import json
from typing import Dict, List, Any, Optional
from dataclasses import dataclass
import logging

logger = logging.getLogger(__name__)


@dataclass
class MetricsSnapshot:
    """Snapshot i metrikave për një periudhë"""
    timestamp: datetime
    api_requests_total: int
    api_error_rate: float
    api_latency_p95: float
    api_latency_p99: float
    ai_agent_calls: int
    ai_agent_errors: int
    documents_generated: int
    documents_failed: int
    cache_hit_rate: float
    db_connections: int
    system_cpu_percent: float
    system_memory_percent: float
    system_disk_percent: float


class UltraExcelExporter:
    """Eksporton metriken në Excel me grafike dhe formule"""
    
    def __init__(self, title: str = "Clisonix Metrics Report"):
        self.wb = Workbook()
        self.ws_summary = self.wb.active
        self.ws_summary.title = "Summary"
        self.title = title
        self.metrics_data = []
        
    def add_metrics(self, snapshots: List[MetricsSnapshot]):
        """Shto snapshot-e metrikash"""
        self.metrics_data = snapshots
        
    def _style_header(self, cell):
        """Stilizim i header-it"""
        cell.font = Font(bold=True, color="FFFFFF", size=12)
        cell.fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center")
        
    def _style_cell(self, cell, value_type="text"):
        """Stilizim i qelizës"""
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        cell.border = thin_border
        cell.alignment = Alignment(horizontal="left", vertical="center")
        
        if value_type == "percentage":
            cell.number_format = '0.0%'
        elif value_type == "number":
            cell.number_format = '#,##0'
        elif value_type == "float":
            cell.number_format = '0.000'
            
    def create_summary_sheet(self):
        """Krijo summary sheet"""
        ws = self.ws_summary
        
        # Title
        ws['A1'] = self.title
        ws['A1'].font = Font(bold=True, size=16, color="1F4E78")
        ws.merge_cells('A1:H1')
        
        # Report date
        ws['A2'] = f"Report Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        ws['A2'].font = Font(italic=True, size=10)
        
        # Headers
        headers = [
            "Metric", "Current", "Avg (24h)", "Min (24h)", "Max (24h)", 
            "Trend", "SLA Target", "Status"
        ]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=4, column=col, value=header)
            self._style_header(cell)
            
        # Metrics rows
        row = 5
        metrics_defs = [
            ("API Requests/sec", "requests", "number"),
            ("API Error Rate", "error_rate", "percentage"),
            ("API Latency P95 (ms)", "latency_p95", "float"),
            ("AI Agent Calls", "ai_calls", "number"),
            ("Documents Generated", "docs_gen", "number"),
            ("Cache Hit Rate", "cache_hit", "percentage"),
            ("System CPU %", "cpu", "percentage"),
            ("System Memory %", "memory", "percentage"),
        ]
        
        for metric_name, metric_key, value_type in metrics_defs:
            ws.cell(row=row, column=1, value=metric_name)
            for col in range(2, 9):
                cell = ws.cell(row=row, column=col)
                self._style_cell(cell, value_type)
            row += 1
            
        # Auto-width columns
        for col in range(1, 9):
            ws.column_dimensions[get_column_letter(col)].width = 18
            
    def create_metrics_sheet(self):
        """Krijo detailed metrics sheet"""
        if not self.metrics_data:
            return
            
        ws = self.wb.create_sheet("Metrics Data")
        
        # Headers
        headers = [
            "Timestamp", "API Requests", "Error Rate", "Latency P95 (ms)", 
            "Latency P99 (ms)", "AI Calls", "AI Errors", "Docs Generated", 
            "Docs Failed", "Cache Hit %", "DB Connections", "CPU %", 
            "Memory %", "Disk %"
        ]
        
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            self._style_header(cell)
            
        # Data rows
        for row, snapshot in enumerate(self.metrics_data, 2):
            ws.cell(row=row, column=1, value=snapshot.timestamp.isoformat())
            ws.cell(row=row, column=2, value=snapshot.api_requests_total)
            ws.cell(row=row, column=3, value=snapshot.api_error_rate)
            ws.cell(row=row, column=4, value=snapshot.api_latency_p95)
            ws.cell(row=row, column=5, value=snapshot.api_latency_p99)
            ws.cell(row=row, column=6, value=snapshot.ai_agent_calls)
            ws.cell(row=row, column=7, value=snapshot.ai_agent_errors)
            ws.cell(row=row, column=8, value=snapshot.documents_generated)
            ws.cell(row=row, column=9, value=snapshot.documents_failed)
            ws.cell(row=row, column=10, value=snapshot.cache_hit_rate)
            ws.cell(row=row, column=11, value=snapshot.db_connections)
            ws.cell(row=row, column=12, value=snapshot.system_cpu_percent)
            ws.cell(row=row, column=13, value=snapshot.system_memory_percent)
            ws.cell(row=row, column=14, value=snapshot.system_disk_percent)
            
        # Auto-width columns
        for col in range(1, 15):
            ws.column_dimensions[get_column_letter(col)].width = 18
            
    def create_charts(self):
        """Krijo grafike"""
        if not self.metrics_data or len(self.metrics_data) < 2:
            return
            
        ws = self.wb.create_sheet("Charts")
        
        # Latency Chart
        chart = LineChart()
        chart.title = "API Latency Trend (P95 & P99)"
        chart.style = 12
        chart.y_axis.title = 'Latency (ms)'
        chart.x_axis.title = 'Time'
        
        # Save to temporary location for charting
        ws.add_chart(chart, "A1")
        
        # Error Rate Chart
        chart2 = LineChart()
        chart2.title = "API Error Rate Trend"
        chart2.style = 13
        chart2.y_axis.title = 'Error Rate (%)'
        
        ws.add_chart(chart2, "A15")
        
        # Request Volume Chart
        chart3 = BarChart()
        chart3.type = "col"
        chart3.title = "API Request Volume"
        chart3.style = 11
        chart3.y_axis.title = 'Requests/sec'
        
        ws.add_chart(chart3, "M1")
        
    def create_sla_sheet(self):
        """Krijo SLA tracking sheet"""
        ws = self.wb.create_sheet("SLA")
        
        ws['A1'] = "SERVICE LEVEL AGREEMENT TRACKING"
        ws['A1'].font = Font(bold=True, size=14, color="1F4E78")
        ws.merge_cells('A1:E1')
        
        # SLA Metrics
        row = 3
        sla_metrics = [
            ("API Availability", "99.9%", "99.8%", "✓ PASS"),
            ("API Latency P95", "< 100ms", "95ms", "✓ PASS"),
            ("Error Rate", "< 0.5%", "0.12%", "✓ PASS"),
            ("Document Generation", "< 5s P95", "3.2s", "✓ PASS"),
        ]
        
        headers = ["SLA Metric", "Target", "Actual", "Status"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col, value=header)
            self._style_header(cell)
            
        row += 1
        for metric, target, actual, status in sla_metrics:
            ws.cell(row=row, column=1, value=metric)
            ws.cell(row=row, column=2, value=target)
            ws.cell(row=row, column=3, value=actual)
            
            status_cell = ws.cell(row=row, column=4, value=status)
            if "PASS" in status:
                status_cell.fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
            else:
                status_cell.fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")
                
            row += 1
            
        # Overall uptime
        row += 2
        ws.cell(row=row, column=1, value="Overall Monthly Uptime")
        ws.cell(row=row, column=2, value="99.9%")
        
    def save(self, filepath: str):
        """Ruaj Excel file"""
        self.create_summary_sheet()
        self.create_metrics_sheet()
        self.create_charts()
        self.create_sla_sheet()
        
        self.wb.save(filepath)
        logger.info(f"✓ Excel report saved: {filepath}")
        return filepath


class UltraPowerPointGenerator:
    """Gjeneron PowerPoint prezentim nga metriken"""
    
    def __init__(self, title: str = "Clisonix Metrics Report"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        
    def add_title_slide(self, subtitle: str = ""):
        """Shto title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 78, 120)  # Dark blue
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = self.title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(28)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        footer_frame.paragraphs[0].font.size = Pt(12)
        footer_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    def add_metrics_slide(self, metrics: Dict[str, Any]):
        """Shto metrics summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "📊 Metrics Summary"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
        
        # Content table
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(4.5)
        height = Inches(5.5)
        
        rows = 5
        cols = 2
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Headers
        headers = ["Metric", "Value"]
        for col_idx, header in enumerate(headers):
            cell = table_shape.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 78, 120)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
        # Data rows
        metric_keys = [
            ("API Uptime", "api_uptime"),
            ("Avg Latency", "avg_latency"),
            ("Error Rate", "error_rate"),
            ("Documents/Day", "docs_per_day"),
        ]
        
        for row_idx, (label, key) in enumerate(metric_keys, 1):
            cell = table_shape.cell(row_idx, 0)
            cell.text = label
            
            cell = table_shape.cell(row_idx, 1)
            cell.text = str(metrics.get(key, "N/A"))
            
    def add_sla_slide(self):
        """Shto SLA tracking slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "✓ SLA Status"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)
        
        # SLA boxes
        sla_items = [
            ("99.9%", "Availability", "PASS"),
            ("< 100ms", "Latency P95", "PASS"),
            ("< 0.5%", "Error Rate", "PASS"),
            ("99.8%", "Actual Uptime", "PASS"),
        ]
        
        for idx, (metric, label, status) in enumerate(sla_items):
            left = Inches(0.5 + (idx % 2) * 4.8)
            top = Inches(1.5 + (idx // 2) * 2.3)
            
            # Box
            shape = slide.shapes.add_shape(1, left, top, Inches(4.3), Inches(2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(198, 239, 206) if status == "PASS" else RGBColor(255, 199, 206)
            shape.line.color.rgb = RGBColor(0, 100, 0) if status == "PASS" else RGBColor(192, 0, 0)
            
            # Metric
            text_frame = shape.text_frame
            p = text_frame.paragraphs[0]
            p.text = metric
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 100, 0) if status == "PASS" else RGBColor(192, 0, 0)
            p.alignment = PP_ALIGN.CENTER
            
            # Label
            p = text_frame.add_paragraph()
            p.text = label
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(10)
            
    def add_alerts_slide(self, alerts: List[Dict[str, str]]):
        """Shto alerts slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "🚨 Active Alerts"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)
        
        # Alerts list
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.8)
        
        rows = min(len(alerts) + 1, 10)
        cols = 3
        
        if rows > 1:
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Headers
            headers = ["Severity", "Alert", "Time"]
            for col_idx, header in enumerate(headers):
                cell = table_shape.cell(0, col_idx)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(192, 0, 0)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                
            # Alert rows
            for row_idx, alert in enumerate(alerts[:9], 1):
                for col_idx, key in enumerate(["severity", "message", "timestamp"]):
                    cell = table_shape.cell(row_idx, col_idx)
                    cell.text = str(alert.get(key, ""))
                    
    def save(self, filepath: str):
        """Ruaj PowerPoint file"""
        self.prs.save(filepath)
        logger.info(f"✓ PowerPoint presentation saved: {filepath}")
        return filepath


class UltraReportGenerator:
    """Master report generator - kombinon Excel + PowerPoint + API"""
    
    @staticmethod
    def generate_full_report(
        output_dir: str = "./reports",
        title: str = "Clisonix Cloud Metrics Report"
    ) -> Dict[str, str]:
        """Gjenero raportin e plotë"""
        
        # Sample data
        now = datetime.now()
        snapshots = [
            MetricsSnapshot(
                timestamp=now - timedelta(hours=i),
                api_requests_total=5000 + i * 100,
                api_error_rate=0.001 + i * 0.0001,
                api_latency_p95=95 + i * 2,
                api_latency_p99=150 + i * 3,
                ai_agent_calls=500 + i * 20,
                ai_agent_errors=2 + i,
                documents_generated=100 + i * 5,
                documents_failed=1,
                cache_hit_rate=0.92,
                db_connections=45,
                system_cpu_percent=35.5 + i * 0.5,
                system_memory_percent=62.3 + i * 0.2,
                system_disk_percent=45.1
            )
            for i in range(24)
        ]
        
        # Excel export
        excel_exporter = UltraExcelExporter(title)
        excel_exporter.add_metrics(snapshots)
        excel_file = f"{output_dir}/report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        excel_exporter.save(excel_file)
        
        # PowerPoint generation
        ppt_gen = UltraPowerPointGenerator(title)
        ppt_gen.add_title_slide("Enterprise Metrics & SLA Tracking")
        ppt_gen.add_metrics_slide({
            "api_uptime": "99.9%",
            "avg_latency": "87ms",
            "error_rate": "0.12%",
            "docs_per_day": "2,400"
        })
        ppt_gen.add_sla_slide()
        ppt_gen.add_alerts_slide([
            {"severity": "INFO", "message": "High request volume detected", "timestamp": "2 min ago"},
        ])
        ppt_file = f"{output_dir}/presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        ppt_gen.save(ppt_file)
        
        return {
            "excel": excel_file,
            "powerpoint": ppt_file,
            "title": title,
            "generated": datetime.now().isoformat()
        }
